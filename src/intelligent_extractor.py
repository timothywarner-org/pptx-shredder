"""
Intelligent PPTX extraction using object model + LLM structural inference.

This module uses the proper PPTX object model to extract structured content,
then calls DeepSeek to infer instructional design patterns and learning structure.
"""

import os
import re
from typing import List, Dict, Any, Optional
from dataclasses import dataclass, asdict
from pathlib import Path

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE
from openai import OpenAI
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

@dataclass
class SlideContent:
    """Structured slide content extracted from PPTX object model."""
    slide_number: int
    title: Optional[str]
    text_content: List[str]
    speaker_notes: str
    bullet_points: List[Dict[str, Any]]  # level, text
    tables: List[Dict[str, Any]]
    images: List[Dict[str, str]]
    charts: List[Dict[str, str]]
    layout_name: str
    slide_size: Dict[str, int]  # width, height

@dataclass 
class InstructionalStructure:
    """LLM-inferred instructional design structure."""
    is_module_start: bool
    module_title: Optional[str] 
    learning_objectives: List[str]
    prerequisites: List[str]
    activity_type: Optional[str]  # lecture, demo, lab, assessment, etc.
    difficulty_level: str  # beginner, intermediate, advanced
    estimated_time_minutes: int
    instructional_notes: str
    content_summary: str

class IntelligentPPTXExtractor:
    """Extract PPTX content using object model + LLM structural inference."""
    
    def __init__(self, pptx_path: str, use_llm: bool = True):
        """Initialize with PPTX file and optional LLM usage."""
        self.pptx_path = Path(pptx_path)
        self.presentation = Presentation(str(self.pptx_path))
        self.use_llm = use_llm
        
        # Initialize DeepSeek client if API key available
        self.client = None
        if use_llm and os.getenv('DEEPSEEK_API_KEY'):
            self.client = OpenAI(
                api_key=os.getenv('DEEPSEEK_API_KEY'),
                base_url="https://api.deepseek.com"
            )
    
    def extract_all_slides(self) -> List[Dict[str, Any]]:
        """Extract content from all slides with LLM structural inference."""
        slides = []
        
        for slide_num, slide in enumerate(self.presentation.slides, 1):
            # Extract structured content using PPTX object model
            slide_content = self._extract_slide_content(slide, slide_num)
            
            # Get LLM structural inference
            if self.use_llm and self.client:
                instructional_structure = self._infer_instructional_structure(slide_content)
            else:
                instructional_structure = self._fallback_structure_detection(slide_content)
            
            # Combine content and structure
            slide_data = {
                'content': asdict(slide_content),
                'structure': asdict(instructional_structure)
            }
            slides.append(slide_data)
        
        return slides
    
    def _extract_slide_content(self, slide, slide_number: int) -> SlideContent:
        """Extract structured content using proper PPTX object model."""
        
        # Extract title using slide layout
        title = None
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()
        
        # Extract all text content with structure preservation
        text_content = []
        bullet_points = []
        
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
                
            if hasattr(shape, 'text_frame') and shape.text_frame:
                # Extract structured text with bullet levels
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        if paragraph.level > 0:  # Bullet point
                            bullet_points.append({
                                'level': paragraph.level,
                                'text': text,
                                'font_size': getattr(paragraph.font, 'size', None),
                                'is_bold': getattr(paragraph.font, 'bold', False)
                            })
                        else:
                            text_content.append(text)
            elif hasattr(shape, 'text') and shape.text.strip():
                text_content.append(shape.text.strip())
        
        # Extract tables with structure
        tables = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_data = self._extract_table_structure(shape)
                if table_data:
                    tables.append(table_data)
        
        # Extract images with metadata
        images = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append({
                    'type': 'image',
                    'position': f"({getattr(shape, 'left', 0)}, {getattr(shape, 'top', 0)})",
                    'size': f"{getattr(shape, 'width', 0)}x{getattr(shape, 'height', 0)}"
                })
        
        # Extract charts
        charts = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                charts.append({
                    'type': 'chart',
                    'position': f"({getattr(shape, 'left', 0)}, {getattr(shape, 'top', 0)})"
                })
        
        # Extract speaker notes
        speaker_notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            speaker_notes = slide.notes_slide.notes_text_frame.text.strip()
        
        # Get slide layout name
        layout_name = slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'Unknown'
        
        return SlideContent(
            slide_number=slide_number,
            title=title,
            text_content=text_content,
            speaker_notes=speaker_notes,
            bullet_points=bullet_points,
            tables=tables,
            images=images,
            charts=charts,
            layout_name=layout_name,
            slide_size={
                'width': self.presentation.slide_width,
                'height': self.presentation.slide_height
            }
        )
    
    def _extract_table_structure(self, shape) -> Optional[Dict[str, Any]]:
        """Extract table structure and sample content."""
        if not hasattr(shape, 'table'):
            return None
        
        table = shape.table
        rows = len(table.rows)
        cols = len(table.columns)
        
        # Extract headers (first row)
        headers = []
        sample_data = []
        
        try:
            if rows > 0:
                first_row = table.rows[0]
                headers = [cell.text.strip() for cell in first_row.cells]
                
                # Extract a few sample rows
                for i in range(1, min(rows, 4)):  # Max 3 sample rows
                    row_data = [cell.text.strip() for cell in table.rows[i].cells]
                    sample_data.append(row_data)
        except:
            pass
        
        return {
            'dimensions': f"{rows}x{cols}",
            'headers': headers,
            'sample_data': sample_data
        }
    
    def _infer_instructional_structure(self, slide_content: SlideContent) -> InstructionalStructure:
        """Use DeepSeek to infer instructional design structure."""
        
        # Prepare content for LLM analysis
        content_summary = self._summarize_slide_for_llm(slide_content)
        
        # Create concise content summary
        bullet_text = ' | '.join([bp['text'][:50] for bp in slide_content.bullet_points[:3]])
        notes_preview = slide_content.speaker_notes[:200] if slide_content.speaker_notes else ""
        
        prompt = f"""Analyze this training slide for instructional design structure.

CONTENT:
Title: {slide_content.title}
Text: {' | '.join(slide_content.text_content[:3])}
Bullets: {bullet_text}
Notes: {notes_preview}

Respond with JSON only:
{{
    "is_module_start": boolean,
    "learning_objectives": ["specific objectives found"],
    "prerequisites": ["prerequisites mentioned"], 
    "activity_type": "lecture|demo|lab|assessment|overview",
    "difficulty_level": "beginner|intermediate|advanced",
    "estimated_time_minutes": number,
    "content_summary": "brief summary"
}}"""

        try:
            response = self.client.chat.completions.create(
                model="deepseek-chat",
                messages=[
                    {"role": "system", "content": "You are an expert in instructional design and technical training. Analyze slides for learning structure and provide structured JSON responses."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.1,
                max_tokens=500
            )
            
            # Parse JSON response
            import json
            response_text = response.choices[0].message.content.strip()
            print(f"DEBUG: DeepSeek response: {response_text[:200]}...")
            
            # Try to extract JSON if response has extra text
            if response_text.startswith('```'):
                # Remove code blocks
                response_text = response_text.split('```')[1]
                if response_text.startswith('json'):
                    response_text = response_text[4:]
            
            result = json.loads(response_text)
            
            return InstructionalStructure(
                is_module_start=result.get('is_module_start', False),
                module_title=slide_content.title if result.get('is_module_start') else None,
                learning_objectives=result.get('learning_objectives', []),
                prerequisites=result.get('prerequisites', []),
                activity_type=result.get('activity_type'),
                difficulty_level=result.get('difficulty_level', 'beginner'),
                estimated_time_minutes=result.get('estimated_time_minutes', 2),
                instructional_notes='LLM analysis completed',
                content_summary=result.get('content_summary', '')
            )
            
        except Exception as e:
            print(f"LLM inference failed for slide {slide_content.slide_number}: {e}")
            return self._fallback_structure_detection(slide_content)
    
    def _fallback_structure_detection(self, slide_content: SlideContent) -> InstructionalStructure:
        """Fallback structure detection without LLM."""
        title = slide_content.title or ""
        
        # Simple heuristics
        is_module_start = any(keyword in title.lower() for keyword in 
                             ['module', 'section', 'chapter', 'overview', 'introduction'])
        
        activity_type = None
        if any(keyword in title.lower() for keyword in ['demo', 'demonstration']):
            activity_type = 'demo'
        elif any(keyword in title.lower() for keyword in ['lab', 'exercise', 'hands-on']):
            activity_type = 'lab'
        elif any(keyword in title.lower() for keyword in ['quiz', 'test', 'assessment']):
            activity_type = 'assessment'
        else:
            activity_type = 'lecture'
        
        return InstructionalStructure(
            is_module_start=is_module_start,
            module_title=title if is_module_start else None,
            learning_objectives=[],
            prerequisites=[],
            activity_type=activity_type,
            difficulty_level='intermediate',
            estimated_time_minutes=3,
            instructional_notes='Fallback detection used',
            content_summary=f"Slide with {len(slide_content.text_content)} text items"
        )
    
    def _summarize_slide_for_llm(self, slide_content: SlideContent) -> str:
        """Create a concise summary of slide content for LLM analysis."""
        summary_parts = []
        
        if slide_content.title:
            summary_parts.append(f"Title: {slide_content.title}")
        
        if slide_content.text_content:
            summary_parts.append(f"Content: {' | '.join(slide_content.text_content[:3])}")
        
        if slide_content.bullet_points:
            bullets = [bp['text'] for bp in slide_content.bullet_points[:3]]
            summary_parts.append(f"Bullets: {' | '.join(bullets)}")
        
        if slide_content.speaker_notes:
            notes_preview = slide_content.speaker_notes[:200] + "..." if len(slide_content.speaker_notes) > 200 else slide_content.speaker_notes
            summary_parts.append(f"Notes: {notes_preview}")
        
        return " | ".join(summary_parts)