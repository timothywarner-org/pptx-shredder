"""
PPTX content extraction module.

Extracts text content, speaker notes, slide titles, and structural information
from PowerPoint presentations while preserving instructional design elements.
"""

import re
from typing import List, Dict, Any, Optional
from pathlib import Path
from dataclasses import dataclass

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass
class SlideData:
    """Container for extracted slide data with comprehensive pedagogical metadata."""
    slide_number: int
    title: Optional[str]
    content: List[str]
    speaker_notes: str
    code_blocks: List[Dict[str, str]]
    is_module_start: bool
    learning_objectives: List[str]
    activity_type: Optional[str]
    # Enhanced pedagogical metadata
    instructor_notes: Dict[str, List[str]]  # categorized by type
    prerequisites: List[str]
    difficulty_level: str
    estimated_time: int  # minutes
    visual_elements: List[Dict[str, str]]
    structured_content: Dict[str, Any]  # preserved formatting
    assessment_items: List[Dict[str, str]]
    compliance_markers: List[str]
    slide_layout_type: str


class PPTXExtractor:
    """Extracts content from PowerPoint presentations with instructional design awareness."""
    
    # Keywords that indicate the start of a new module/section (enterprise training focused)
    MODULE_MARKERS = [
        'module', 'section', 'chapter', 'unit', 'lesson', 'part', 'topic', 'agenda', 'overview',
        'phase', 'step', 'stage', 'milestone', 'objective', 'demo', 'workshop', 'session',
        'lab', 'exercise', 'hands-on', 'practice', 'activity', 'scenario', 'case study',
        'walkthrough', 'tutorial', 'deep dive', 'checkpoint'
    ]
    
    # Numbered module patterns
    MODULE_NUMBER_PATTERNS = [
        r'(?:module|section|chapter|unit|lesson|part|topic|phase|step)\s*\d+',
        r'\d+\.\s*(?:module|section|chapter|unit|lesson|part|topic|phase|step)',
        r'(?:step|phase)\s+[ivx]+',  # Roman numerals
        r'^\d+[\.\)]\s+'  # Simple numbered items
    ]
    
    # Keywords that indicate learning activities (enterprise training focused)
    ACTIVITY_MARKERS = {
        'lab': 'hands-on-lab',
        'exercise': 'guided-exercise', 
        'practice': 'practice-session',
        'demo': 'demonstration',
        'demonstration': 'demonstration',
        'try it': 'hands-on-activity',
        'hands-on': 'hands-on-activity',
        'activity': 'learning-activity',
        'assignment': 'assignment',
        'quiz': 'knowledge-check',
        'test': 'assessment',
        'assessment': 'formal-assessment',
        'review': 'knowledge-review',
        'troubleshooting': 'troubleshooting-scenario',
        'case study': 'case-study',
        'scenario': 'scenario-based-learning',
        'best practice': 'best-practices',
        'real world': 'real-world-application',
        'certification': 'certification-prep'
    }
    
    # Instructor note categorization patterns
    INSTRUCTOR_NOTE_PATTERNS = {
        'timing': [r'(?:time|duration|minutes?):', r'(?:spend|allow|take)\s+\d+\s*(?:min|minutes?)'],
        'emphasis': [r'(?:important|key|critical|note|remember):', r'(?:emphasize|highlight|stress)'],
        'examples': [r'(?:example|instance|case|scenario):', r'for example', r'such as'],
        'tips': [r'(?:tip|hint|suggestion):', r'pro tip', r'best practice'],
        'warnings': [r'(?:warning|caution|avoid|don\'t):', r'be careful', r'watch out'],
        'context': [r'(?:context|background|why):', r'the reason', r'this is because'],
        'delivery': [r'(?:say|tell|explain|mention):', r'make sure to', r'don\'t forget']
    }
    
    # Difficulty indicators
    DIFFICULTY_MARKERS = {
        'beginner': ['basic', 'introduction', 'fundamentals', 'getting started', 'overview'],
        'intermediate': ['intermediate', 'advanced concepts', 'diving deeper', 'detailed'],
        'advanced': ['advanced', 'expert', 'deep dive', 'complex', 'enterprise', 'production']
    }
    
    # Compliance and certification markers
    COMPLIANCE_MARKERS = [
        'certification', 'certified', 'compliance', 'audit', 'security', 
        'gdpr', 'hipaa', 'sox', 'iso', 'nist', 'pci', 'regulation'
    ]
    
    def __init__(self, pptx_path: str):
        """Initialize extractor with PowerPoint file path."""
        self.pptx_path = Path(pptx_path)
        self.presentation = Presentation(str(self.pptx_path))
        
    def extract(self) -> List[SlideData]:
        """Extract all content from the presentation."""
        slides_data = []
        
        for slide_num, slide in enumerate(self.presentation.slides, 1):
            slide_data = self._extract_slide(slide, slide_num)
            slides_data.append(slide_data)
            
        return slides_data
    
    def _extract_slide(self, slide, slide_number: int) -> SlideData:
        """Extract content from a single slide with robust error handling."""
        try:
            title = self._extract_title(slide)
            content = self._extract_content(slide)
            speaker_notes = self._extract_speaker_notes(slide)
            
            # Validate minimum content
            if not title and not content and not speaker_notes:
                title = f"Slide {slide_number}"
                content = ["[Slide content could not be extracted]"]
            
            # Extract with individual error handling
            try:
                code_blocks = self._extract_code_blocks(slide)
            except Exception:
                code_blocks = []
            
            try:
                is_module_start = self._is_module_start(title, content)
            except Exception:
                is_module_start = False
            
            try:
                learning_objectives = self._extract_learning_objectives(content, speaker_notes)
            except Exception:
                learning_objectives = []
            
            try:
                activity_type = self._detect_activity_type(title, content)
            except Exception:
                activity_type = None
            
            try:
                instructor_notes = self._categorize_instructor_notes(speaker_notes)
            except Exception:
                instructor_notes = {}
            
            try:
                prerequisites = self._extract_prerequisites(content, speaker_notes)
            except Exception:
                prerequisites = []
            
            try:
                difficulty_level = self._assess_difficulty_level(title, content, speaker_notes)
            except Exception:
                difficulty_level = 'beginner'
            
            try:
                estimated_time = self._estimate_slide_time(content, speaker_notes, activity_type)
            except Exception:
                estimated_time = 2  # Default 2 minutes
            
            try:
                visual_elements = self._extract_visual_elements(slide)
            except Exception:
                visual_elements = []
            
            try:
                structured_content = self._extract_structured_content(slide)
            except Exception:
                structured_content = {'lists': [], 'emphasized_text': [], 'headings': [], 'layout_sections': []}
            
            try:
                assessment_items = self._extract_assessment_items(content, speaker_notes)
            except Exception:
                assessment_items = []
            
            try:
                compliance_markers = self._extract_compliance_markers(content, speaker_notes)
            except Exception:
                compliance_markers = []
            
            try:
                slide_layout_type = self._detect_slide_layout(slide)
            except Exception:
                slide_layout_type = 'standard-content'
                
        except Exception as e:
            # Critical error - return minimal viable slide data
            return SlideData(
                slide_number=slide_number,
                title=f"Slide {slide_number} (Extraction Failed)",
                content=["[Content extraction failed]"],
                speaker_notes="",
                code_blocks=[],
                is_module_start=False,
                learning_objectives=[],
                activity_type=None,
                instructor_notes={},
                prerequisites=[],
                difficulty_level='beginner',
                estimated_time=2,
                visual_elements=[],
                structured_content={'lists': [], 'emphasized_text': [], 'headings': [], 'layout_sections': []},
                assessment_items=[],
                compliance_markers=[],
                slide_layout_type='standard-content'
            )
        
        return SlideData(
            slide_number=slide_number,
            title=title,
            content=content,
            speaker_notes=speaker_notes,
            code_blocks=code_blocks,
            is_module_start=is_module_start,
            learning_objectives=learning_objectives,
            activity_type=activity_type,
            instructor_notes=instructor_notes,
            prerequisites=prerequisites,
            difficulty_level=difficulty_level,
            estimated_time=estimated_time,
            visual_elements=visual_elements,
            structured_content=structured_content,
            assessment_items=assessment_items,
            compliance_markers=compliance_markers,
            slide_layout_type=slide_layout_type
        )
    
    def _extract_title(self, slide) -> Optional[str]:
        """Extract slide title using shape type and position heuristics."""
        # Try to get title from slide layout first
        if slide.shapes.title:
            return slide.shapes.title.text.strip()
        
        # Fallback: look for text at the top of the slide
        title_candidates = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                # Check if shape is positioned like a title (top 25% of slide)
                if hasattr(shape, 'top') and shape.top < self.presentation.slide_height * 0.25:
                    title_candidates.append((shape.text.strip(), shape.top))
        
        if title_candidates:
            # Return the topmost text as title
            title_candidates.sort(key=lambda x: x[1])
            return title_candidates[0][0]
        
        return None
    
    def _extract_content(self, slide) -> List[str]:
        """Extract all text content from slide shapes."""
        content = []
        
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue  # Skip title, handled separately
                
            if hasattr(shape, 'text') and shape.text.strip():
                content.append(shape.text.strip())
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Handle grouped shapes
                content.extend(self._extract_from_group(shape))
            elif hasattr(shape, 'text_frame'):
                # Handle text frames
                text = self._extract_from_text_frame(shape.text_frame)
                if text:
                    content.append(text)
        
        return content
    
    def _extract_from_group(self, group_shape) -> List[str]:
        """Extract text from grouped shapes."""
        content = []
        for shape in group_shape.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                content.append(shape.text.strip())
        return content
    
    def _extract_from_text_frame(self, text_frame) -> str:
        """Extract formatted text from text frame."""
        if not text_frame.text.strip():
            return ""
        
        # For now, just return plain text
        # TODO: Preserve formatting for lists, emphasis, etc.
        return text_frame.text.strip()
    
    def _extract_speaker_notes(self, slide) -> str:
        """Extract speaker notes from slide."""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                return notes_slide.notes_text_frame.text.strip()
        return ""
    
    def _extract_code_blocks(self, slide) -> List[Dict[str, str]]:
        """Identify and extract code blocks from slide content."""
        code_blocks = []
        
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                text = shape.text.strip()
                
                # Heuristics for identifying code:
                # 1. Monospace font indicators
                # 2. Common code patterns
                # 3. Indentation patterns
                if self._looks_like_code(text, shape):
                    language = self._detect_language(text)
                    code_blocks.append({
                        'code': text,
                        'language': language
                    })
        
        return code_blocks
    
    def _looks_like_code(self, text: str, shape) -> bool:
        """Determine if text looks like code using various heuristics."""
        # Check for common code indicators
        code_indicators = [
            '{', '}', '()', '[]', ';', '->', '=>', 
            'function', 'def ', 'class ', 'import ', 'from ',
            'SELECT', 'INSERT', 'UPDATE', 'DELETE',
            '$', '#', '//', '/*', '*/', '<!--', '-->'
        ]
        
        text_lower = text.lower()
        indicator_count = sum(1 for indicator in code_indicators if indicator in text_lower)
        
        # If multiple indicators present, likely code
        if indicator_count >= 2:
            return True
        
        # Check for indentation patterns (common in code)
        lines = text.split('\n')
        indented_lines = sum(1 for line in lines if line.startswith(('  ', '\t')))
        if len(lines) > 1 and indented_lines / len(lines) > 0.3:
            return True
        
        return False
    
    def _detect_language(self, code: str) -> str:
        """Attempt to detect programming language from code content."""
        code_lower = code.lower()
        
        # Simple language detection based on keywords (order matters - more specific first)
        if any(keyword in code_lower for keyword in ['select ', 'insert ', 'update ', 'delete ', 'where ']):
            return 'sql'
        elif any(keyword in code_lower for keyword in ['def ', 'import ', 'from ', 'print(']):
            return 'python'
        elif any(keyword in code_lower for keyword in ['function', 'var ', 'let ', 'const ', 'console.log']):
            return 'javascript'
        elif any(keyword in code_lower for keyword in ['<div', '<span', '<html', '<body']):
            return 'html'
        elif any(keyword in code_lower for keyword in ['public class', 'private ', 'public static']):
            return 'java'
        elif any(keyword in code_lower for keyword in ['using ', 'namespace', 'public class']):
            return 'csharp'
        
        return 'text'  # Default fallback
    
    def _is_module_start(self, title: Optional[str], content: List[str]) -> bool:
        """Enhanced module detection with fuzzy matching and content analysis."""
        if not title:
            return False
        
        title_lower = title.lower().strip()
        
        # Check for explicit module markers
        if any(marker in title_lower for marker in self.MODULE_MARKERS):
            return True
        
        # Check for numbered patterns
        for pattern in self.MODULE_NUMBER_PATTERNS:
            if re.search(pattern, title_lower):
                return True
        
        # Check content for module indicators
        all_text = ' '.join(content).lower()
        module_content_indicators = [
            'learning objectives', 'what you will learn', 'in this module', 'module overview',
            'objectives:', 'goals:', 'by the end', 'after completing', 'you will be able',
            'agenda', 'outline', 'topics covered'
        ]
        
        if any(indicator in all_text for indicator in module_content_indicators):
            return True
        
        # Special case: if slide has very little content and seems like a section divider
        if len(title_lower.split()) <= 5 and len(all_text) < 100:
            # Check for section-like patterns
            section_patterns = [
                r'part\s+\d+', r'section\s+\d+', r'chapter\s+\d+',
                r'introduction', r'getting started', r'overview',
                r'conclusion', r'summary', r'wrap.?up'
            ]
            if any(re.search(pattern, title_lower) for pattern in section_patterns):
                return True
        
        return False
    
    def _extract_learning_objectives(self, content: List[str], speaker_notes: str) -> List[str]:
        """Extract learning objectives with robust, inclusive patterns."""
        objectives = []
        all_text = ' '.join(content) + ' ' + speaker_notes
        
        # Focused, high-quality patterns for learning objectives
        objective_patterns = [
            # Explicit objective statements with proper context
            r'(?:learning\s+objectives?|objectives?|goals?)[:\s]*\n?[•\-\*]?\s*([A-Z][^.\n!?]{15,120})',
            
            # "You will" patterns with action verbs
            r'(?:you|students?|learners?|participants?)\s+(?:will|can|should)\s+(?:be\s+able\s+to\s+)?(learn|understand|identify|demonstrate|explain|configure|implement|analyze|create|evaluate|apply|assess|manage|administer|deploy|troubleshoot|validate|enable)\s+([^.\n!?]{10,100})',
            
            # "After this" clear completion patterns
            r'(?:by\s+the\s+end\s+of\s+this|after\s+completing\s+this|upon\s+completion)[^.\n!?]*?(?:you|students?|learners?)\s+(?:will|should)\s+(?:be\s+able\s+to\s+)?([^.\n!?]{15,100})',
            
            # Bullet point objectives with action verbs
            r'[•\-\*]\s*(?:Be\s+able\s+to\s+|Learn\s+to\s+|Understand\s+how\s+to\s+)?([A-Z][a-z]+\s+(?:GHAS|GitHub|security|policies|features|access|requirements)[^•\-\*\n]{10,80})',
            
            # Clear instructional outcomes
            r'(?:learning\s+outcomes?|outcomes?)[:\s]*[•\-\*]\s*([A-Z][^•\-\*\n]{15,100})'
        ]
        
        for pattern in objective_patterns:
            matches = re.finditer(pattern, all_text, re.IGNORECASE | re.MULTILINE)
            for match in matches:
                objective = match.group(1).strip() if match.lastindex else match.group(0).strip()
                
                # More lenient filtering
                if (8 <= len(objective) <= 150 and 
                    not objective.lower().startswith(('ing ', '. ', 'the ', 'and ', 'or ')) and
                    objective.count(' ') >= 1):  # At least 2 words
                    
                    # Clean up common artifacts
                    objective = re.sub(r'^(?:and\s+|or\s+)', '', objective, flags=re.IGNORECASE)
                    objective = objective.strip('.,!?')
                    
                    if len(objective) > 5:
                        objectives.append(objective)
        
        # Remove duplicates while preserving order
        seen = set()
        unique_objectives = []
        for obj in objectives:
            obj_clean = obj.lower().strip('.,!? ')
            if obj_clean not in seen and len(obj_clean) > 5:
                seen.add(obj_clean)
                unique_objectives.append(obj)
        
        return unique_objectives[:8]  # Allow more objectives to be captured
    
    def _detect_activity_type(self, title: Optional[str], content: List[str]) -> Optional[str]:
        """Detect the type of learning activity represented by the slide."""
        if not title:
            return None
        
        title_lower = title.lower()
        all_content = ' '.join(content).lower()
        
        # Check title and content for activity markers
        for marker, activity_type in self.ACTIVITY_MARKERS.items():
            if marker in title_lower or marker in all_content:
                return activity_type
        
        return None
    
    def _categorize_instructor_notes(self, speaker_notes: str) -> Dict[str, List[str]]:
        """Categorize speaker notes by pedagogical intent."""
        categories = {
            'timing': [],
            'emphasis': [],
            'examples': [],
            'tips': [],
            'warnings': [],
            'context': [],
            'delivery': []
        }
        
        if not speaker_notes:
            return categories
            
        # Split notes into sentences for analysis
        sentences = re.split(r'[.!?]+', speaker_notes)
        
        for sentence in sentences:
            sentence = sentence.strip()
            if len(sentence) < 5:  # Skip very short fragments
                continue
                
            # Check each category pattern
            for category, patterns in self.INSTRUCTOR_NOTE_PATTERNS.items():
                for pattern in patterns:
                    if re.search(pattern, sentence, re.IGNORECASE):
                        categories[category].append(sentence)
                        break
        
        # Remove empty categories
        return {k: v for k, v in categories.items() if v}
    
    def _extract_prerequisites(self, content: List[str], speaker_notes: str) -> List[str]:
        """Extract prerequisites with robust, inclusive patterns."""
        prerequisites = []
        all_text = ' '.join(content) + ' ' + speaker_notes
        
        # Simplified, more inclusive patterns
        prereq_patterns = [
            # Direct requirements (no punctuation requirement)
            r'(?:requires?|needs?|must\s+have|should\s+have)\s+([^.\n!?]{5,80})',
            r'(?:prerequisite|requirement)[s]?[:\s]*([^.\n!?]{5,80})',
            
            # Experience/Knowledge patterns
            r'(?:experience|familiarity|knowledge)\s+(?:with|of|in)\s+([^.\n!?]{5,80})',
            r'(?:prior|previous)\s+(?:experience|knowledge|familiarity)\s+(?:with|of|in)\s+([^.\n!?]{5,80})',
            
            # Access patterns (common in enterprise)
            r'(?:admin|administrator|administrative)\s+(?:access|permissions?|rights?)',
            r'(?:access\s+to|permissions?\s+(?:for|to))\s+([^.\n!?]{5,80})',
            
            # License patterns
            r'(?:licens[es]*|subscription)\s+(?:for|to|of)\s+([^.\n!?]{5,80})',
            r'(?:GHAS|GitHub\s+Advanced\s+Security)\s+licens[es]*',
            
            # Basic/fundamental requirements
            r'(?:basic|fundamental|working)\s+(?:understanding|knowledge|familiarity)\s+(?:of|with)\s+([^.\n!?]{5,80})',
            
            # Before starting patterns
            r'before\s+(?:starting|beginning|taking)[^.\n!?]*?(?:you|students?)\s+(?:should|must|need)[^.\n!?]*?([^.\n!?]{8,80})',
            
            # Assumes patterns
            r'(?:assumes?|assuming)\s+(?:you\s+have\s+|that\s+you\s+have\s+|)?([^.\n!?]{8,80})',
            
            # Simple bullet patterns
            r'[•\-\*]\s*([^•\-\*\n]*(?:license|access|permission|experience|knowledge|understanding|familiarity)[^•\-\*\n]*)',
        ]
        
        for pattern in prereq_patterns:
            matches = re.finditer(pattern, all_text, re.IGNORECASE | re.MULTILINE)
            for match in matches:
                if match.lastindex and match.lastindex >= 1:
                    prereq = match.group(1).strip()
                else:
                    prereq = match.group(0).strip()
                
                # Clean up common prefixes and suffixes
                prereq = re.sub(r'^(?:prerequisite|requirement)[s]?[:\-]?\s*', '', prereq, flags=re.IGNORECASE)
                prereq = re.sub(r'^(?:you\s+(?:should\s+)?have\s+)', '', prereq, flags=re.IGNORECASE)
                prereq = prereq.strip('.,!? ')
                
                # More lenient validation
                if (5 <= len(prereq) <= 100 and 
                    not prereq.lower().startswith(('for ', 'in ', 'of ', 'to ', 'and ', 'or ', 'the ')) and
                    prereq.count(' ') >= 0):  # Allow single words for licenses, etc.
                    
                    prerequisites.append(prereq)
        
        # Remove duplicates and clean up
        seen = set()
        unique_prereqs = []
        for prereq in prerequisites:
            prereq_clean = prereq.lower().strip()
            if prereq_clean not in seen and len(prereq_clean) > 3:
                seen.add(prereq_clean)
                unique_prereqs.append(prereq)
        
        return unique_prereqs[:5]  # Allow more prerequisites
    
    def _assess_difficulty_level(self, title: Optional[str], content: List[str], speaker_notes: str) -> str:
        """Assess the difficulty level of the slide content."""
        all_text = ' '.join([title or ''] + content + [speaker_notes]).lower()
        
        # Score each difficulty level
        scores = {'beginner': 0, 'intermediate': 0, 'advanced': 0}
        
        for level, markers in self.DIFFICULTY_MARKERS.items():
            for marker in markers:
                scores[level] += all_text.count(marker)
        
        # Determine difficulty based on scores and content complexity
        max_score = max(scores.values())
        if max_score == 0:
            # Use heuristics: code blocks and technical terms suggest higher difficulty
            code_indicators = len(re.findall(r'[{}();]', all_text))
            if code_indicators > 5:
                return 'advanced'
            elif code_indicators > 2:
                return 'intermediate'
            else:
                return 'beginner'
        
        return max(scores, key=scores.get)
    
    def _estimate_slide_time(self, content: List[str], speaker_notes: str, activity_type: Optional[str]) -> int:
        """Estimate time needed for slide in minutes."""
        # Base time calculation
        content_length = sum(len(text) for text in content)
        notes_length = len(speaker_notes)
        
        # Base time: ~150 words per minute reading speed
        base_time = (content_length + notes_length) / (150 * 5)  # Rough chars per word
        
        # Activity type multipliers
        activity_multipliers = {
            'hands-on-lab': 10,
            'guided-exercise': 5,
            'practice-session': 3,
            'demonstration': 2,
            'hands-on-activity': 4,
            'troubleshooting-scenario': 8,
            'case-study': 6
        }
        
        multiplier = activity_multipliers.get(activity_type, 1)
        estimated_minutes = max(1, int(base_time * multiplier))
        
        return min(estimated_minutes, 45)  # Cap at 45 minutes
    
    def _extract_visual_elements(self, slide) -> List[Dict[str, str]]:
        """Extract detailed information about visual elements on the slide."""
        visual_elements = []
        
        for shape in slide.shapes:
            element = {}
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                element = {
                    'type': 'image',
                    'description': self._describe_image_context(shape, slide),
                    'position': f'top={getattr(shape, "top", 0)}, left={getattr(shape, "left", 0)}',
                    'size': f'width={getattr(shape, "width", 0)}, height={getattr(shape, "height", 0)}'
                }
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_info = self._extract_table_info(shape)
                element = {
                    'type': 'table',
                    'description': table_info['description'],
                    'data': table_info['summary'],
                    'structure': table_info['structure']
                }
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart_info = self._extract_chart_info(shape)
                element = {
                    'type': 'chart',
                    'description': chart_info['description'],
                    'chart_type': chart_info['chart_type'],
                    'data_summary': chart_info['data_summary']
                }
            elif shape.shape_type == MSO_SHAPE_TYPE.DIAGRAM:
                element = {
                    'type': 'diagram',
                    'description': 'SmartArt diagram or flowchart',
                    'content': 'Process flow or conceptual diagram'
                }
            elif hasattr(shape, 'text') and shape.text.strip():
                # Text boxes with special formatting
                text_content = shape.text.strip()
                if len(text_content) < 100 and any(keyword in text_content.lower() for keyword in ['screenshot', 'figure', 'diagram', 'example', 'demo']):
                    element = {
                        'type': 'caption',
                        'description': f'Text caption: {text_content}',
                        'content': text_content
                    }
            
            if element:
                visual_elements.append(element)
        
        return visual_elements
    
    def _describe_image_context(self, shape, slide) -> str:
        """Generate contextual description for images based on surrounding content."""
        # Look for nearby text that might describe the image
        surrounding_text = []
        
        if hasattr(shape, 'top') and hasattr(shape, 'left'):
            shape_top = getattr(shape, 'top', 0)
            shape_left = getattr(shape, 'left', 0)
            
            # Find text shapes near this image
            for other_shape in slide.shapes:
                if (hasattr(other_shape, 'text') and other_shape.text.strip() and 
                    hasattr(other_shape, 'top') and hasattr(other_shape, 'left')):
                    
                    other_top = getattr(other_shape, 'top', 0)
                    other_left = getattr(other_shape, 'left', 0)
                    
                    # Check if text is near the image (rough proximity)
                    if (abs(other_top - shape_top) < 100000 or  # Near vertically
                        abs(other_left - shape_left) < 100000):   # Near horizontally
                        text = other_shape.text.strip()
                        if len(text) < 200 and any(keyword in text.lower() for keyword in ['screenshot', 'figure', 'example', 'diagram', 'ui', 'interface', 'demo', 'console', 'terminal']):
                            surrounding_text.append(text)
        
        if surrounding_text:
            return f"Image with context: {' | '.join(surrounding_text[:2])}"
        else:
            # Generic description based on slide context
            slide_title = self._extract_title(slide)
            if slide_title:
                title_lower = slide_title.lower()
                if 'github' in title_lower:
                    return "GitHub interface screenshot or diagram"
                elif any(term in title_lower for term in ['command', 'cli', 'terminal']):
                    return "Command line interface or terminal screenshot"
                elif 'demo' in title_lower:
                    return "Demonstration screenshot or workflow diagram"
                elif any(term in title_lower for term in ['overview', 'architecture', 'structure']):
                    return "System architecture or overview diagram"
            
            return "Technical diagram, screenshot, or instructional image"
    
    def _extract_table_info(self, shape) -> Dict[str, str]:
        """Extract meaningful information from table shapes."""
        if not hasattr(shape, 'table'):
            return {
                'description': 'Data table',
                'summary': 'Table structure not accessible',
                'structure': 'Unknown dimensions'
            }
        
        table = shape.table
        rows = len(table.rows) if hasattr(table, 'rows') else 0
        cols = len(table.columns) if hasattr(table, 'columns') else 0
        
        # Try to extract some sample content
        sample_content = []
        try:
            if rows > 0 and cols > 0:
                # Get first row (likely headers)
                first_row_text = []
                for cell in table.rows[0].cells:
                    if hasattr(cell, 'text') and cell.text.strip():
                        first_row_text.append(cell.text.strip())
                
                if first_row_text:
                    sample_content = first_row_text[:3]  # First 3 headers
        except:
            pass
        
        description = f"Data table with {rows} rows and {cols} columns"
        if sample_content:
            description += f" (columns: {', '.join(sample_content)})"
        
        return {
            'description': description,
            'summary': f"Tabular data with {rows}x{cols} structure",
            'structure': f"{rows} rows × {cols} columns"
        }
    
    def _extract_chart_info(self, shape) -> Dict[str, str]:
        """Extract information about chart elements."""
        # Basic chart type detection would require more complex analysis
        # For now, provide contextual description
        
        chart_type = "data visualization"
        description = "Chart or graph showing data relationships"
        data_summary = "Quantitative data presentation"
        
        # Could be enhanced with chart.xml analysis in future
        return {
            'description': description,
            'chart_type': chart_type,
            'data_summary': data_summary
        }
    
    def _extract_structured_content(self, slide) -> Dict[str, Any]:
        """Extract content with preserved structure and formatting."""
        structured = {
            'lists': [],
            'emphasized_text': [],
            'headings': [],
            'layout_sections': []
        }
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                # Extract list structures
                if hasattr(shape.text_frame, 'paragraphs'):
                    list_items = []
                    for para in shape.text_frame.paragraphs:
                        if para.level > 0:  # Indented = list item
                            list_items.append({
                                'text': para.text.strip(),
                                'level': para.level,
                                'bullet_style': 'bullet'  # Could be enhanced to detect actual style
                            })
                    if list_items:
                        structured['lists'].append(list_items)
                
                # Extract emphasized text (would need run-level analysis for bold/italic)
                text = shape.text.strip()
                if text and len(text) < 100:  # Likely emphasized if short
                    # Simple heuristic: ALL CAPS suggests emphasis
                    if text.isupper() and len(text) > 3:
                        structured['emphasized_text'].append(text)
        
        return structured
    
    def _extract_assessment_items(self, content: List[str], speaker_notes: str) -> List[Dict[str, str]]:
        """Extract quiz questions, assessments, and knowledge checks."""
        assessment_items = []
        all_text = ' '.join(content) + ' ' + speaker_notes
        
        # Question patterns
        question_patterns = [
            r'(What (?:is|are|do|does)[^?]*\\?)',
            r'(How (?:do|does|can|will)[^?]*\\?)',
            r'(Why (?:is|are|do|does)[^?]*\\?)',
            r'(Which (?:of|one)[^?]*\\?)',
            r'(True or False[^?]*\\?)'
        ]
        
        for pattern in question_patterns:
            matches = re.finditer(pattern, all_text, re.IGNORECASE)
            for match in matches:
                question = match.group(1).strip()
                assessment_items.append({
                    'type': 'question',
                    'content': question,
                    'format': 'multiple_choice' if 'which' in question.lower() else 'open_ended'
                })
        
        return assessment_items
    
    def _extract_compliance_markers(self, content: List[str], speaker_notes: str) -> List[str]:
        """Extract compliance and certification related markers."""
        compliance_items = []
        all_text = ' '.join(content + [speaker_notes]).lower()
        
        for marker in self.COMPLIANCE_MARKERS:
            if marker in all_text:
                compliance_items.append(marker.upper())
        
        return list(set(compliance_items))  # Remove duplicates
    
    def _detect_slide_layout(self, slide) -> str:
        """Detect the semantic layout type of the slide."""
        shape_types = [shape.shape_type for shape in slide.shapes]
        
        # Simple layout detection based on content
        if MSO_SHAPE_TYPE.TABLE in shape_types:
            return 'data-table'
        elif MSO_SHAPE_TYPE.CHART in shape_types:
            return 'data-visualization'
        elif MSO_SHAPE_TYPE.PICTURE in shape_types:
            return 'image-focused'
        elif len([s for s in slide.shapes if hasattr(s, 'text') and len(s.text) > 100]) > 2:
            return 'content-heavy'
        elif slide.shapes.title and len(slide.shapes) <= 3:
            return 'title-slide'
        else:
            return 'standard-content'