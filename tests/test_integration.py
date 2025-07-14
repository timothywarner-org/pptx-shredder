"""
Integration tests for the complete PPTX Shredder pipeline.
"""

import pytest
import tempfile
import shutil
from pathlib import Path
from click.testing import CliRunner

from src.shred import shred
from src.extractor import PPTXExtractor
from src.formatter import MarkdownFormatter


class TestIntegration:
    """Integration tests for the full pipeline."""
    
    def test_complete_pipeline_with_sample_pptx(self, sample_pptx, temp_dir):
        """Test the complete pipeline from PPTX to markdown."""
        output_dir = temp_dir / "output"
        
        # Extract
        extractor = PPTXExtractor(str(sample_pptx))
        slides_data = extractor.extract()
        
        # Format
        formatter = MarkdownFormatter()
        markdown_files = formatter.format(slides_data, sample_pptx.stem)
        
        # Write files
        output_dir.mkdir()
        for filename, content in markdown_files.items():
            output_file = output_dir / filename
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(content)
        
        # Verify files were created
        created_files = list(output_dir.glob("*.md"))
        assert len(created_files) >= 1
        
        # Check file content
        for md_file in created_files:
            content = md_file.read_text(encoding='utf-8')
            assert content.startswith("---\n")  # YAML frontmatter
            assert "# " in content  # At least one heading
            assert "Azure" in content or "Cloud" in content  # Expected content
    
    def test_cli_with_specific_file(self, sample_pptx, temp_dir):
        """Test CLI with a specific PPTX file."""
        runner = CliRunner()
        output_dir = temp_dir / "cli_output"
        
        result = runner.invoke(shred, [
            str(sample_pptx),
            '--output-dir', str(output_dir),
            '--verbose'
        ])
        
        assert result.exit_code == 0
        assert "Processing Complete" in result.output
        
        # Check that output files were created
        created_files = list(output_dir.glob("*.md"))
        assert len(created_files) >= 1
    
    def test_cli_dry_run(self, sample_pptx):
        """Test CLI dry run mode."""
        runner = CliRunner()
        
        result = runner.invoke(shred, [
            str(sample_pptx),
            '--dry-run'
        ])
        
        assert result.exit_code == 0
        assert "Dry run complete" in result.output
        assert "test_presentation.pptx" in result.output
    
    def test_cli_with_input_directory(self, sample_pptx, temp_dir):
        """Test CLI scanning input directory."""
        # Create input directory and copy sample file
        input_dir = temp_dir / "input_test"
        output_dir = temp_dir / "output_test"
        input_dir.mkdir()
        
        # Copy sample PPTX to input directory
        test_file = input_dir / "presentation.pptx"
        shutil.copy2(sample_pptx, test_file)
        
        runner = CliRunner()
        result = runner.invoke(shred, [
            '--input-dir', str(input_dir),
            '--output-dir', str(output_dir)
        ])
        
        assert result.exit_code == 0
        assert "presentation.pptx" in result.output
        
        # Check output files
        created_files = list(output_dir.glob("*.md"))
        assert len(created_files) >= 1
    
    def test_cli_no_files_found(self, temp_dir):
        """Test CLI behavior when no PPTX files are found."""
        empty_input_dir = temp_dir / "empty_input"
        empty_input_dir.mkdir()
        
        runner = CliRunner()
        result = runner.invoke(shred, [
            '--input-dir', str(empty_input_dir)
        ])
        
        assert result.exit_code == 0
        assert "No PPTX files found" in result.output
    
    def test_cli_invalid_file(self, temp_dir):
        """Test CLI behavior with invalid file."""
        invalid_file = temp_dir / "not_a_pptx.txt"
        invalid_file.write_text("This is not a PowerPoint file")
        
        runner = CliRunner()
        result = runner.invoke(shred, [str(invalid_file)])
        
        # Should complete but warn about skipping non-PPTX file
        assert result.exit_code == 0
        assert "No PPTX files found" in result.output or "Skipping" in result.output
    
    @pytest.mark.parametrize("strategy", ['instructional', 'sequential', 'module-based'])
    def test_cli_different_strategies(self, sample_pptx, temp_dir, strategy):
        """Test CLI with different chunking strategies."""
        output_dir = temp_dir / f"output_{strategy}"
        
        runner = CliRunner()
        result = runner.invoke(shred, [
            str(sample_pptx),
            '--output-dir', str(output_dir),
            '--strategy', strategy
        ])
        
        assert result.exit_code == 0
        
        # Check that output files were created
        created_files = list(output_dir.glob("*.md"))
        assert len(created_files) >= 1
    
    def test_cli_custom_chunk_size(self, sample_pptx, temp_dir):
        """Test CLI with custom chunk size."""
        output_dir = temp_dir / "output_custom"
        
        runner = CliRunner()
        result = runner.invoke(shred, [
            str(sample_pptx),
            '--output-dir', str(output_dir),
            '--chunk-size', '500'  # Small size to force more chunks
        ])
        
        assert result.exit_code == 0
        
        # Should create output files
        created_files = list(output_dir.glob("*.md"))
        assert len(created_files) >= 1
    
    def test_error_handling_corrupted_pptx(self, temp_dir):
        """Test error handling with corrupted PPTX file."""
        # Create a fake PPTX file (just text with .pptx extension)
        fake_pptx = temp_dir / "corrupted.pptx"
        fake_pptx.write_text("This is not a real PowerPoint file")
        
        runner = CliRunner()
        result = runner.invoke(shred, [str(fake_pptx)])
        
        # Should handle error gracefully
        assert result.exit_code == 0  # CLI should not crash
        assert "❌ corrupted.pptx failed" in result.output or "No PPTX files found" in result.output
    
    def test_output_file_content_quality(self, sample_pptx, temp_dir):
        """Test that output files contain expected high-quality content."""
        output_dir = temp_dir / "quality_test"
        
        runner = CliRunner()
        result = runner.invoke(shred, [
            str(sample_pptx),
            '--output-dir', str(output_dir),
            '--verbose'
        ])
        
        assert result.exit_code == 0
        
        # Read and verify content quality
        created_files = list(output_dir.glob("*.md"))
        assert len(created_files) >= 1
        
        for md_file in created_files:
            content = md_file.read_text(encoding='utf-8')
            
            # Check YAML frontmatter structure
            assert content.startswith("---\n")
            yaml_end = content.find("\n---\n", 4)
            assert yaml_end > 0
            
            frontmatter = content[4:yaml_end]
            assert "module_id:" in frontmatter
            assert "module_title:" in frontmatter
            assert "slide_range:" in frontmatter
            
            # Check markdown content
            markdown_content = content[yaml_end + 5:]
            assert "# " in markdown_content  # Has headings
            assert "Azure" in markdown_content or "Cloud" in markdown_content  # Has expected content
            assert "**Instructor Notes:**" in markdown_content  # Has speaker notes
    
    @pytest.mark.slow
    def test_performance_with_large_presentation(self, temp_dir):
        """Test performance with a larger presentation."""
        # Create a larger test presentation
        from pptx import Presentation
        
        large_pptx = temp_dir / "large_presentation.pptx"
        prs = Presentation()
        
        # Create 50 slides
        for i in range(50):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = f"Slide {i+1}: Content Title"
            if slide.placeholders:
                slide.placeholders[1].text = f"This is content for slide {i+1}. " * 10
        
        prs.save(str(large_pptx))
        
        output_dir = temp_dir / "large_output"
        
        import time
        start_time = time.time()
        
        runner = CliRunner()
        result = runner.invoke(shred, [
            str(large_pptx),
            '--output-dir', str(output_dir)
        ])
        
        elapsed_time = time.time() - start_time
        
        assert result.exit_code == 0
        assert elapsed_time < 30  # Should complete within 30 seconds
        
        # Should create multiple chunks due to size
        created_files = list(output_dir.glob("*.md"))
        assert len(created_files) >= 1