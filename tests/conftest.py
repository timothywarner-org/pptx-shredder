"""
Pytest configuration and shared fixtures.
"""

import pytest
import tempfile
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from src.extractor import SlideData


@pytest.fixture
def temp_dir():
    """Create a temporary directory for test files."""
    temp_path = Path(tempfile.mkdtemp())
    yield temp_path
    shutil.rmtree(temp_path)


@pytest.fixture
def sample_pptx(temp_dir):
    """Create a sample PPTX file for testing."""
    pptx_path = temp_dir / "test_presentation.pptx"
    
    # Create presentation
    prs = Presentation()
    
    # Slide 1: Module start
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Module 1: Azure Fundamentals"
    slide.placeholders[1].text = "Introduction to Cloud Computing"
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Learning objective: Students will understand cloud computing basics."
    
    # Slide 2: Content slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "What is Cloud Computing?"
    slide.placeholders[1].text = """• On-demand self-service
• Broad network access
• Resource pooling
• Rapid elasticity
• Measured service"""
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Explain each characteristic with examples."
    
    # Slide 3: Lab slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Lab: Create Azure Account"
    slide.placeholders[1].text = """1. Navigate to portal.azure.com
2. Click 'Create account'
3. Follow registration steps
4. Verify email address"""
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Allow 15 minutes for account creation."
    
    # Save the presentation
    prs.save(str(pptx_path))
    return pptx_path


@pytest.fixture
def sample_slide_data():
    """Create sample SlideData objects for testing."""
    return [
        SlideData(
            slide_number=1,
            title="Module 1: Azure Fundamentals", 
            content=["Introduction to Cloud Computing"],
            speaker_notes="Learning objective: Students will understand cloud computing basics.",
            code_blocks=[],
            is_module_start=True,
            learning_objectives=["Students will understand cloud computing basics"],
            activity_type=None
        ),
        SlideData(
            slide_number=2,
            title="What is Cloud Computing?",
            content=["• On-demand self-service", "• Broad network access", "• Resource pooling"],
            speaker_notes="Explain each characteristic with examples.",
            code_blocks=[],
            is_module_start=False,
            learning_objectives=[],
            activity_type=None
        ),
        SlideData(
            slide_number=3,
            title="Lab: Create Azure Account",
            content=["1. Navigate to portal.azure.com", "2. Click 'Create account'"],
            speaker_notes="Allow 15 minutes for account creation.",
            code_blocks=[],
            is_module_start=False,
            learning_objectives=[],
            activity_type="lab"
        )
    ]


@pytest.fixture
def sample_code_slide_data():
    """Create sample SlideData with code blocks for testing."""
    return SlideData(
        slide_number=4,
        title="Python Example",
        content=["Here's a simple Python function:"],
        speaker_notes="Walk through the code step by step.",
        code_blocks=[
            {
                'code': 'def hello_world():\n    print("Hello, World!")\n    return True',
                'language': 'python'
            }
        ],
        is_module_start=False,
        learning_objectives=[],
        activity_type=None
    )